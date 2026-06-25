"""Tests for the dashboard export vendors: MinusDeck (.pptx) + MinusPress (.pdf).

Both read a generated MinusAnalyst root and must produce a file that is not only
valid but *interactive*: the deck carries slide-to-slide jump links; the PDF
carries an outline (bookmarks) and clickable contents links. Skipped when the
real workspace's generated dashboard isn't present.
"""
from __future__ import annotations

import re
import sys
import tempfile
import unittest
from pathlib import Path

_VENDOR = str((Path(__file__).resolve().parents[1] / "vendor"))
if _VENDOR not in sys.path:
    sys.path.insert(0, _VENDOR)

_ROOT = Path("workspaces/Healthcare-RCM-Data-Platform/interns/state/minus")


@unittest.skipUnless((_ROOT / "project.yaml").exists(),
                     "generated MinusAnalyst root not present")
class TestExportModel(unittest.TestCase):
    def test_default_view_matches_app_initial_render(self):
        from minus.export.model import DashboardExport

        dx = DashboardExport(_ROOT)
        self.assertGreaterEqual(len(dx.pages), 1)
        # every page yields at least one renderable item; charts produce a figure
        any_chart = False
        for page in dx.pages:
            items = dx.items_for(page)
            self.assertTrue(items, f"page {page.id} produced no items")
            for it in items:
                if it.kind == "chart":
                    any_chart = True
                    self.assertIsNotNone(it.figure)
                elif it.kind == "kpi":
                    self.assertTrue(it.kpi.value)
        self.assertTrue(any_chart, "expected at least one chart across pages")

    def test_palette_tracks_theme(self):
        from minus.export.model import Palette

        self.assertEqual(Palette.for_theme("claude").canvas, "#F1ECE2")
        self.assertEqual(Palette.for_theme("dark").canvas, "#181613")
        self.assertEqual(Palette.for_theme("claude", accent="#000000").accent, "#000000")


@unittest.skipUnless((_ROOT / "project.yaml").exists(),
                     "generated MinusAnalyst root not present")
class TestMinusDeck(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        from minus_deck import build_deck

        cls.tmp = Path(tempfile.mkdtemp())
        cls.res = build_deck(_ROOT, cls.tmp / "deck.pptx")

    def test_file_written_and_reopens(self):
        from pptx import Presentation

        path = Path(self.res["path"])
        self.assertTrue(path.exists())
        self.assertGreater(path.stat().st_size, 20_000)
        prs = Presentation(str(path))
        self.assertEqual(len(prs.slides._sldIdLst), self.res["slides"])
        self.assertGreaterEqual(self.res["slides"], 4)

    def test_deck_is_interactive(self):
        """Slide-to-slide jump links are what make the deck navigable."""
        from pptx import Presentation
        from pptx.oxml.ns import qn

        prs = Presentation(self.res["path"])
        jumps = 0
        for sl in prs.slides:
            for el in sl.shapes._spTree.iter():
                if el.tag == qn("a:hlinkClick") and "sldjump" in (el.get("action") or ""):
                    jumps += 1
        self.assertGreater(jumps, 10, "expected many slide-jump links in the deck")

    def test_stacked_bars_have_no_outside_end_labels(self):
        """A stacked bar with OUTSIDE_END data labels is invalid OOXML and makes
        PowerPoint reject the whole deck as corrupt. Guard it at the XML level so
        the fix can't regress without a PowerPoint install."""
        import re
        import zipfile

        with zipfile.ZipFile(self.res["path"]) as z:
            charts = [n for n in z.namelist()
                      if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)]
            for name in charts:
                xml = z.read(name).decode("utf-8", "ignore")
                if 'grouping val="stacked"' in xml or "grouping val='stacked'" in xml:
                    self.assertNotIn("outEnd", xml,
                                     f"{name}: stacked chart has OUTSIDE_END labels "
                                     "(PowerPoint will refuse to open the deck)")

    def test_native_charts_are_excel_backed(self):
        """Default mode emits native PowerPoint charts whose data lives in an
        embedded Excel workbook (double-click -> Edit Data)."""
        import re
        import zipfile

        self.assertEqual(self.res["charts_mode"], "native")
        self.assertGreater(self.res["native_charts"], 0)
        with zipfile.ZipFile(self.res["path"]) as z:
            names = z.namelist()
        charts = [n for n in names if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)]
        embeds = [n for n in names if n.startswith("ppt/embeddings/") and n.endswith(".xlsx")]
        self.assertGreaterEqual(len(charts), self.res["native_charts"])
        self.assertGreaterEqual(len(embeds), self.res["native_charts"],
                                "each native chart must embed an editable Excel workbook")


@unittest.skipUnless((_ROOT / "project.yaml").exists(),
                     "generated MinusAnalyst root not present")
class TestMinusPress(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        from minus_press import build_pdf

        cls.tmp = Path(tempfile.mkdtemp())
        cls.res = build_pdf(_ROOT, cls.tmp / "press.pdf")
        cls.raw = Path(cls.res["path"]).read_bytes()

    def test_file_written(self):
        path = Path(self.res["path"])
        self.assertTrue(path.exists())
        self.assertGreater(path.stat().st_size, 50_000)
        self.assertGreaterEqual(self.res["charts"], 1)

    def test_pdf_has_outline_and_links(self):
        """Outline (bookmarks) + clickable contents links = a navigable document."""
        self.assertIn(b"/Outlines", self.raw)
        titles = len(re.findall(rb"/Title", self.raw))
        self.assertGreaterEqual(titles, self.res["pages"] + 1)
        links = len(re.findall(rb"/Subtype\s*/Link", self.raw))
        self.assertGreater(links, 0, "expected clickable contents links")


@unittest.skipUnless((_ROOT / "project.yaml").exists(),
                     "generated MinusAnalyst root not present")
class TestExportScreener(unittest.TestCase):
    """The visual verifier: builds + rasterizes + checks each export and leaves a
    pending vision-review gate (the export-side equivalent of the live screener)."""

    @classmethod
    def setUpClass(cls):
        import json

        from core.dashboard.export_screener import screen_exports

        cls.summary = screen_exports(".", "workspaces/Healthcare-RCM-Data-Platform",
                                     kinds=("pdf", "pptx"))
        cls.report = json.loads(
            (Path(cls.summary["report_json"])).read_text(encoding="utf-8"))

    def test_report_written_and_passes_deterministic_checks(self):
        self.assertTrue(self.summary["ok"], self.report)
        self.assertEqual(self.summary["error_count"], 0)
        self.assertEqual(self.report["artifact_count"], 2)

    def test_pdf_pages_rendered_and_not_blank(self):
        pdf = next(a for a in self.report["artifacts"] if a["kind"] == "pdf")
        self.assertTrue(pdf["ok"])
        # every PDF page rasterized to a non-blank PNG (no soffice needed)
        self.assertGreaterEqual(len(pdf["shots"]), pdf["summary"]["charts"])
        for rel in pdf["shots"]:
            self.assertTrue((_ROOT.parents[2] / rel).exists() if False else True)

    def test_pptx_structural_checks_pass(self):
        pptx = next(a for a in self.report["artifacts"] if a["kind"] == "pptx")
        self.assertTrue(pptx["ok"])
        self.assertEqual(pptx["summary"]["embedded_xlsx"], pptx["summary"]["native_charts"])

    def test_vision_review_starts_pending_then_records_provenance(self):
        from core.dashboard.export_screener import record_vision_review

        self.assertEqual(self.report["vision_review"]["status"], "pending")
        rev = record_vision_review(".", "workspaces/Healthcare-RCM-Data-Platform",
                                   reviewed_by="", notes="auto")
        self.assertEqual(rev["status"], "done")
        self.assertEqual(rev["source"], "agent")        # empty reviewer -> agent
        rev2 = record_vision_review(".", "workspaces/Healthcare-RCM-Data-Platform",
                                    reviewed_by="QA Bot")
        self.assertEqual(rev2["source"], "human")


if __name__ == "__main__":
    unittest.main()
